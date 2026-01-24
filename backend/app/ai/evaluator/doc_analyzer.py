from __future__ import annotations

import re
from difflib import SequenceMatcher
from pathlib import Path
from typing import Any


def calculate_readability(text: str) -> dict[str, float]:
    """
    Calculate readability metrics (Flesch Reading Ease and Grade Level).
    Returns scores between 0-100 (higher = easier to read).
    """
    if not text or not text.strip():
        return {"flesch_reading_ease": 0.0, "flesch_kincaid_grade": 0.0}

    # Count sentences
    sentences = re.split(r'[.!?]+', text)
    sentences = [s.strip() for s in sentences if s.strip()]
    num_sentences = max(len(sentences), 1)

    # Count words
    words = re.findall(r'\b\w+\b', text)
    num_words = max(len(words), 1)

    # Count syllables (approximation)
    def count_syllables(word: str) -> int:
        word = word.lower()
        vowels = "aeiouy"
        syllable_count = 0
        previous_was_vowel = False

        for char in word:
            is_vowel = char in vowels
            if is_vowel and not previous_was_vowel:
                syllable_count += 1
            previous_was_vowel = is_vowel

        # Adjust for silent 'e'
        if word.endswith('e') and syllable_count > 1:
            syllable_count -= 1

        return max(syllable_count, 1)

    total_syllables = sum(count_syllables(w) for w in words)

    # Flesch Reading Ease: 206.835 - 1.015 * (words/sentences) - 84.6 * (syllables/words)
    words_per_sentence = num_words / num_sentences
    syllables_per_word = total_syllables / num_words

    flesch_reading_ease = 206.835 - 1.015 * words_per_sentence - 84.6 * syllables_per_word
    flesch_reading_ease = max(0.0, min(100.0, flesch_reading_ease))

    # Flesch-Kincaid Grade Level: 0.39 * (words/sentences) + 11.8 * (syllables/words) - 15.59
    flesch_kincaid_grade = 0.39 * words_per_sentence + 11.8 * syllables_per_word - 15.59
    flesch_kincaid_grade = max(0.0, flesch_kincaid_grade)

    return {
        "flesch_reading_ease": round(flesch_reading_ease, 2),
        "flesch_kincaid_grade": round(flesch_kincaid_grade, 2),
    }


def check_basic_plagiarism(text: str, reference_texts: list[str] | None = None) -> dict[str, Any]:
    """
    Basic plagiarism detection using sequence matching.
    Returns similarity percentage with reference texts.
    """
    if not reference_texts or not text.strip():
        return {
            "plagiarism_detected": False,
            "max_similarity": 0.0,
            "similar_sources": [],
        }

    text_lower = text.lower()
    similarities = []

    for idx, ref in enumerate(reference_texts):
        if not ref or not ref.strip():
            continue

        ref_lower = ref.lower()
        similarity = SequenceMatcher(None, text_lower, ref_lower).ratio()

        if similarity > 0.3:  # Threshold for suspicion
            similarities.append({
                "source_index": idx,
                "similarity": round(similarity * 100, 2),
            })

    max_similarity = max((s["similarity"] for s in similarities), default=0.0)
    plagiarism_detected = max_similarity > 70.0  # 70% similarity threshold

    return {
        "plagiarism_detected": plagiarism_detected,
        "max_similarity": round(max_similarity, 2),
        "similar_sources": sorted(similarities, key=lambda x: x["similarity"], reverse=True)[:3],
    }


def analyze_structure(text: str) -> dict[str, Any]:
    """
    Analyze document structure quality.
    Checks for paragraphs, headings, formatting, etc.
    """
    if not text or not text.strip():
        return {"structure_quality": 0.0, "has_paragraphs": False, "paragraph_count": 0}

    paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
    paragraph_count = len(paragraphs)

    # Check for headings (lines that start with # or are all caps)
    lines = text.split('\n')
    potential_headings = 0
    for line in lines:
        line = line.strip()
        if line.startswith('#') or (line.isupper() and len(line.split()) <= 10 and len(line) > 0):
            potential_headings += 1

    # Calculate structure quality score (0-100)
    quality_score = 0.0

    # Has multiple paragraphs
    if paragraph_count >= 2:
        quality_score += 30
    elif paragraph_count == 1:
        quality_score += 10

    # Has headings/sections
    if potential_headings >= 3:
        quality_score += 30
    elif potential_headings >= 1:
        quality_score += 15

    # Paragraph length consistency (not too short, not too long)
    if paragraphs:
        avg_para_length = sum(len(p.split()) for p in paragraphs) / len(paragraphs)
        if 30 <= avg_para_length <= 150:
            quality_score += 20
        elif 15 <= avg_para_length <= 200:
            quality_score += 10

    # Has proper spacing
    if '\n\n' in text:
        quality_score += 20

    quality_score = min(100.0, quality_score)

    return {
        "structure_quality": round(quality_score, 2),
        "has_paragraphs": paragraph_count > 1,
        "paragraph_count": paragraph_count,
        "heading_count": potential_headings,
        "avg_paragraph_words": round(sum(len(p.split()) for p in paragraphs) / max(len(paragraphs), 1), 2),
    }


def analyze_text(
    *,
    text: str,
    keywords: list[str] | None = None,
    min_words: int = 0,
    reference_texts: list[str] | None = None,
    enable_readability: bool = True,
    enable_plagiarism: bool = False,
    enable_structure: bool = True,
) -> dict[str, Any]:
    """
    Comprehensive text analysis with multiple metrics.

    Args:
        text: The text to analyze
        keywords: List of keywords to search for
        min_words: Minimum expected word count
        reference_texts: Reference texts for plagiarism detection
        enable_readability: Calculate readability metrics
        enable_plagiarism: Check for plagiarism
        enable_structure: Analyze document structure

    Returns:
        Dictionary with analysis results including:
        - word_count
        - keywords_found
        - readability_score (if enabled)
        - plagiarism_detected (if enabled)
        - structure_quality (if enabled)
    """

    text = text or ""

    # Basic word count and keyword matching
    words = re.findall(r'\b\w+\b', text)
    word_count = len(words)

    keyword_list = [k.strip() for k in (keywords or []) if isinstance(k, str) and k.strip()]
    lower_text = text.lower()
    found_keywords = []

    for keyword in keyword_list:
        # Support multi-word keywords
        if keyword.lower() in lower_text:
            found_keywords.append(keyword)

    result = {
        "word_count": word_count,
        "keywords_found": found_keywords,
        "keyword_match_ratio": round(len(found_keywords) / max(len(keyword_list), 1) * 100, 2),
        "meets_min_words": word_count >= min_words,
    }

    # Readability analysis
    if enable_readability:
        readability = calculate_readability(text)
        result["readability_score"] = readability.get("flesch_reading_ease", 0.0)
        result["grade_level"] = readability.get("flesch_kincaid_grade", 0.0)

    # Plagiarism check
    if enable_plagiarism and reference_texts:
        plagiarism = check_basic_plagiarism(text, reference_texts)
        result["plagiarism_detected"] = plagiarism.get("plagiarism_detected", False)
        result["max_similarity"] = plagiarism.get("max_similarity", 0.0)
        result["similar_sources"] = plagiarism.get("similar_sources", [])

    # Structure analysis
    if enable_structure:
        structure = analyze_structure(text)
        result["structure_quality"] = structure.get("structure_quality", 0.0)
        result["paragraph_count"] = structure.get("paragraph_count", 0)
        result["has_paragraphs"] = structure.get("has_paragraphs", False)

    return result


def extract_text_from_pdf(path: str | Path) -> str:
    """Extract text content from a PDF file."""
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"PDF file not found: {p}")

    try:
        from pypdf import PdfReader  # type: ignore
    except Exception:
        try:
            from PyPDF2 import PdfReader  # type: ignore
        except Exception as e:
            raise RuntimeError("PDF extraction library not installed (pypdf or PyPDF2 required)") from e

    try:
        reader = PdfReader(str(p))
        chunks: list[str] = []

        for page in reader.pages:
            text = page.extract_text() or ""
            if text.strip():
                chunks.append(text)

        return "\n\n".join(chunks)
    except Exception as e:
        raise RuntimeError(f"Failed to extract text from PDF: {e}") from e


def extract_text_from_pptx(path: str | Path) -> str:
    """Extract text content from a PowerPoint PPTX file."""
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"PPTX file not found: {p}")

    try:
        from pptx import Presentation  # type: ignore
    except Exception as e:
        raise RuntimeError("PPTX extraction library not installed (python-pptx required)") from e

    try:
        prs = Presentation(str(p))
        chunks: list[str] = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)

                # Extract text from table cells if present
                if hasattr(shape, "table"):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text:
                                slide_text.append(cell.text)

            if slide_text:
                chunks.append(f"Slide {slide_num}:\n" + "\n".join(slide_text))

        return "\n\n".join(chunks)
    except Exception as e:
        raise RuntimeError(f"Failed to extract text from PPTX: {e}") from e


def extract_text_from_docx(path: str | Path) -> str:
    """Extract text content from a Word DOCX file."""
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"DOCX file not found: {p}")

    try:
        from docx import Document  # type: ignore
    except Exception as e:
        raise RuntimeError("DOCX extraction library not installed (python-docx required)") from e

    try:
        doc = Document(str(p))
        chunks: list[str] = []

        # Extract text from paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                chunks.append(paragraph.text)

        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    chunks.append(" | ".join(row_text))

        return "\n\n".join(chunks)
    except Exception as e:
        raise RuntimeError(f"Failed to extract text from DOCX: {e}") from e


def extract_text_from_doc(path: str | Path) -> str:
    """
    Extract text content from an old Word DOC file.
    Note: This requires antiword or similar external tool on the system.
    Falls back to basic extraction if antiword is not available.
    """
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"DOC file not found: {p}")

    import subprocess

    # Try antiword first (works on Linux/Mac)
    try:
        result = subprocess.run(
            ["antiword", str(p)],
            capture_output=True,
            text=True,
            timeout=10,
            check=False
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout.strip()
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass

    # Fallback: Try python-docx (might work for some .doc files)
    try:
        return extract_text_from_docx(p)
    except Exception:
        pass

    raise RuntimeError(
        "Failed to extract text from DOC file. "
        "Please install 'antiword' or convert to DOCX format."
    )


def extract_text_from_attachment(path: str | Path) -> str:
    """
    Auto-detect file type and extract text from common document formats.
    Supports: PDF, PPTX, DOCX, DOC, TXT.

    Args:
        path: Path to the document file

    Returns:
        Extracted text content

    Raises:
        FileNotFoundError: If file doesn't exist
        RuntimeError: If extraction fails or format is unsupported
    """
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"File not found: {p}")

    ext = p.suffix.lower()

    if ext == ".pdf":
        return extract_text_from_pdf(p)
    elif ext == ".pptx" or ext == ".ppt":
        return extract_text_from_pptx(p)
    elif ext == ".docx":
        return extract_text_from_docx(p)
    elif ext == ".doc":
        return extract_text_from_doc(p)
    elif ext == ".txt":
        return p.read_text(encoding="utf-8", errors="ignore")
    else:
        raise RuntimeError(f"Unsupported file type: {ext}")


async def analyze_text_with_groq(
    *,
    user_uid: str,
    text: str,
    keywords: list[str] | None = None,
    min_words: int = 0,
    task_title: str = "",
    task_description: str = "",
    reference_texts: list[str] | None = None,
    enable_readability: bool = True,
    enable_plagiarism: bool = False,
    enable_structure: bool = True,
) -> dict[str, Any]:
    """
    Comprehensive text analysis with Groq AI enhancement.

    First performs rule-based analysis, then enhances with Groq insights.

    Args:
        user_uid: Student's user ID (for rate limiting)
        text: The text to analyze
        keywords: List of keywords to search for
        min_words: Minimum expected word count
        task_title: Assignment title
        task_description: Assignment description
        reference_texts: Reference texts for plagiarism detection
        enable_readability: Calculate readability metrics
        enable_plagiarism: Check for plagiarism
        enable_structure: Analyze document structure

    Returns:
        Dictionary with analysis results including Groq insights
    """
    import logging
    logger = logging.getLogger(__name__)

    # First, run rule-based analysis
    base_result = analyze_text(
        text=text,
        keywords=keywords,
        min_words=min_words,
        reference_texts=reference_texts,
        enable_readability=enable_readability,
        enable_plagiarism=enable_plagiarism,
        enable_structure=enable_structure,
    )

    # Try to enhance with Groq
    try:
        from app.services.groq_service import groq_service

        if not groq_service.is_available():
            logger.debug("Groq not available, using rule-based analysis only")
            return base_result

        # Prepare data for Groq
        found_keywords = base_result.get("keywords_found", [])
        keyword_list = [k.strip() for k in (keywords or []) if isinstance(k, str) and k.strip()]
        missing_keywords = [k for k in keyword_list if k not in found_keywords]

        readability = {
            "flesch_score": base_result.get("readability_score", 0),
            "grade_level": base_result.get("grade_level", 0)
        }

        # Call Groq for analysis
        groq_result = await groq_service.analyze_document(
            user_uid=user_uid,
            content=text,
            word_count=base_result.get("word_count", 0),
            required_keywords=keyword_list,
            found_keywords=found_keywords,
            missing_keywords=missing_keywords,
            readability=readability,
            task_title=task_title,
            task_description=task_description,
            min_words=min_words
        )

        # Merge Groq results with base results
        base_result["groq_analysis"] = {
            "quality_assessment": groq_result.get("quality_assessment", ""),
            "structure_feedback": groq_result.get("structure_feedback", ""),
            "improvements": groq_result.get("improvements", []),
            "suggested_score": groq_result.get("suggested_score"),
            "raw_response": groq_result.get("raw_response", "")
        }

        # Use Groq's suggested score if available
        if groq_result.get("suggested_score") is not None:
            base_result["groq_suggested_score"] = groq_result["suggested_score"]

        return base_result

    except Exception as e:
        logger.warning(f"Failed to get Groq analysis: {e}")
        return base_result


def preprocess_document_for_groq(
    text: str,
    keywords: list[str] | None,
    base_metrics: dict[str, Any],
    task_title: str = "",
    task_description: str = "",
    min_words: int = 0
) -> dict[str, Any]:
    """
    Preprocess document data for Groq analysis.

    This extracts and structures the relevant information
    to minimize the Groq prompt size.

    Args:
        text: The submission text
        keywords: Required keywords
        base_metrics: Results from rule-based analyze_text()
        task_title: Assignment title
        task_description: Assignment description
        min_words: Minimum word requirement

    Returns:
        Preprocessed data dict ready for Groq
    """
    # Truncate content preview
    content_preview = text[:3000] if len(text) > 3000 else text

    # Process keywords
    keyword_list = [k.strip() for k in (keywords or []) if isinstance(k, str) and k.strip()]
    found_keywords = base_metrics.get("keywords_found", [])
    missing_keywords = [k for k in keyword_list if k not in found_keywords]

    return {
        "content_preview": content_preview,
        "word_count": base_metrics.get("word_count", 0),
        "required_keywords": keyword_list,
        "found_keywords": found_keywords,
        "missing_keywords": missing_keywords,
        "readability": {
            "flesch_score": base_metrics.get("readability_score", 0),
            "grade_level": base_metrics.get("grade_level", 0)
        },
        "task_requirements": {
            "title": task_title,
            "description": task_description[:500] if task_description else "",
            "min_words": min_words
        }
    }
